"""
PPTX Tool - Microsoft PowerPoint presentation processing.

Provides capabilities for:
- Reading presentations
- Extracting text and notes
- Slide enumeration
- Basic structure analysis
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from tools.base import Tool, ToolParameter, ToolResult, register_tool

logger = logging.getLogger("titan.tools.documents.pptx")


class PPTXTool(Tool):
    """
    Microsoft PowerPoint presentation tool.

    Allows agents to read and analyze PPTX files.
    """

    @property
    def name(self) -> str:
        return "pptx"

    @property
    def description(self) -> str:
        return (
            "Microsoft PowerPoint (PPTX) presentation tool. "
            "Actions: 'outline' (get presentation structure), 'read' (extract all text), "
            "'slide' (get specific slide), 'notes' (get speaker notes)."
        )

    @property
    def parameters(self) -> list[ToolParameter]:
        return [
            ToolParameter(
                name="action",
                type="string",
                description="Action: 'outline', 'read', 'slide', 'notes'",
                required=True,
                enum=["outline", "read", "slide", "notes"],
            ),
            ToolParameter(
                name="path",
                type="string",
                description="Path to PPTX file",
                required=True,
            ),
            ToolParameter(
                name="slide_number",
                type="integer",
                description="Slide number (1-indexed) for 'slide' action",
                required=False,
            ),
            ToolParameter(
                name="max_chars",
                type="integer",
                description="Maximum characters to return (default: 50000)",
                required=False,
            ),
        ]

    async def execute(self, **kwargs: Any) -> ToolResult:
        try:
            from pptx import Presentation
        except ImportError:
            return ToolResult(
                success=False,
                output=None,
                error="python-pptx is required. Install with: pip install 'agentic-titan[documents]'",
            )

        action = kwargs.get("action", "")
        path = kwargs.get("path", "")
        max_chars = kwargs.get("max_chars", 50000)

        if not path:
            return ToolResult(
                success=False,
                output=None,
                error="'path' parameter is required",
            )

        try:
            prs = Presentation(path)

            if action == "outline":
                outline = []
                for i, slide in enumerate(prs.slides, start=1):
                    slide_info = {
                        "slide_number": i,
                        "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
                        "shapes": len(slide.shapes),
                    }

                    # Get title if present
                    if slide.shapes.title:
                        slide_info["title"] = slide.shapes.title.text

                    # Check for notes
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text.strip():
                            slide_info["has_notes"] = True

                    outline.append(slide_info)

                return ToolResult(
                    success=True,
                    output={
                        "slide_count": len(prs.slides),
                        "outline": outline,
                    },
                )

            elif action == "read":
                all_text = []

                for i, slide in enumerate(prs.slides, start=1):
                    slide_text = [f"[Slide {i}]"]

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)

                    all_text.append("\n".join(slide_text))

                text = "\n\n".join(all_text)
                if len(text) > max_chars:
                    text = text[:max_chars] + "\n\n[Truncated]"

                return ToolResult(
                    success=True,
                    output={
                        "text": text,
                        "slide_count": len(prs.slides),
                        "total_chars": sum(len(t) for t in all_text),
                    },
                )

            elif action == "slide":
                slide_number = kwargs.get("slide_number")
                if not slide_number:
                    return ToolResult(
                        success=False,
                        output=None,
                        error="'slide_number' parameter required for slide action",
                    )

                if slide_number < 1 or slide_number > len(prs.slides):
                    return ToolResult(
                        success=False,
                        output=None,
                        error=f"Invalid slide number {slide_number}. Presentation has {len(prs.slides)} slides.",
                    )

                slide = prs.slides[slide_number - 1]

                # Extract all text content
                content = []
                title = None

                for shape in slide.shapes:
                    if shape == slide.shapes.title:
                        title = shape.text if hasattr(shape, "text") else None
                    elif hasattr(shape, "text") and shape.text.strip():
                        content.append({
                            "type": shape.shape_type.name if hasattr(shape, "shape_type") else "Unknown",
                            "text": shape.text,
                        })

                # Get notes if present
                notes = None
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text

                return ToolResult(
                    success=True,
                    output={
                        "slide_number": slide_number,
                        "title": title,
                        "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
                        "content": content,
                        "notes": notes,
                        "shape_count": len(slide.shapes),
                    },
                )

            elif action == "notes":
                notes_data = []

                for i, slide in enumerate(prs.slides, start=1):
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        if notes_text.strip():
                            notes_data.append({
                                "slide_number": i,
                                "notes": notes_text,
                            })

                return ToolResult(
                    success=True,
                    output={
                        "slide_count": len(prs.slides),
                        "slides_with_notes": len(notes_data),
                        "notes": notes_data,
                    },
                )

            else:
                return ToolResult(
                    success=False,
                    output=None,
                    error=f"Unknown action: {action}",
                )

        except FileNotFoundError:
            return ToolResult(
                success=False,
                output=None,
                error=f"File not found: {path}",
            )
        except Exception as e:
            logger.exception(f"PPTX tool error: {e}")
            return ToolResult(
                success=False,
                output=None,
                error=f"Failed to process PPTX: {e}",
            )


# Register the tool
pptx_tool = PPTXTool()
register_tool(pptx_tool)
